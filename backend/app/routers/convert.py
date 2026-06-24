import io
import zipfile
from pathlib import Path

import fitz
from docx import Document
from fastapi import APIRouter, File, Form, HTTPException, UploadFile
from fastapi.responses import StreamingResponse
from openpyxl import Workbook
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from app.constants import (
    IMAGE_TO_PDF_EXTS,
    MIME_DOCX,
    MIME_HTML,
    MIME_PDF,
    MIME_PPTX,
    MIME_TXT,
    MIME_XLSX,
    MIME_ZIP,
    HTML_TO_PDF_EXTS,
    OFFICE_TO_PDF_EXTS,
    TO_PDF_EXTS,
)
from app.converters.html import HtmlConversionError, html_to_pdf
from app.converters.office import OfficeConversionError, convert_to_pdf
from app.validation import read_validated_pdf, read_validated_upload

router = APIRouter(prefix="/api/convert", tags=["convert"])


def _pdf_to_docx_bytes(pdf_bytes: bytes) -> bytes:
    doc = Document()
    pdf = fitz.open(stream=pdf_bytes, filetype="pdf")

    for page_index in range(len(pdf)):
        page = pdf[page_index]
        text = page.get_text("text").strip()
        if text:
            doc.add_paragraph(text)
        else:
            doc.add_paragraph(f"[Page {page_index + 1} — no extractable text]")
        if page_index < len(pdf) - 1:
            doc.add_page_break()

    pdf.close()
    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer.read()


def _pdf_to_xlsx_bytes(pdf_bytes: bytes) -> bytes:
    workbook = Workbook()
    workbook.remove(workbook.active)
    pdf = fitz.open(stream=pdf_bytes, filetype="pdf")

    for page_index in range(len(pdf)):
        page = pdf[page_index]
        sheet = workbook.create_sheet(title=f"Page {page_index + 1}"[:31])
        tables = page.find_tables()

        if not tables.tables:
            sheet.cell(row=1, column=1, value=page.get_text("text").strip() or f"[Page {page_index + 1}]")
            continue

        row_offset = 1
        for table in tables:
            for row in table.extract():
                for col_index, value in enumerate(row, start=1):
                    sheet.cell(row=row_offset, column=col_index, value=value)
                row_offset += 1
            row_offset += 1

    pdf.close()
    buffer = io.BytesIO()
    workbook.save(buffer)
    buffer.seek(0)
    return buffer.read()


def _pdf_to_pptx_bytes(pdf_bytes: bytes) -> bytes:
    presentation = Presentation()
    blank_layout = presentation.slide_layouts[6]
    pdf = fitz.open(stream=pdf_bytes, filetype="pdf")

    for page_index in range(len(pdf)):
        page = pdf[page_index]
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        image_bytes = pix.tobytes("png")

        slide = presentation.slides.add_slide(blank_layout)
        image_stream = io.BytesIO(image_bytes)
        slide.shapes.add_picture(
            image_stream,
            Inches(0),
            Inches(0),
            width=presentation.slide_width,
            height=presentation.slide_height,
        )

    pdf.close()
    buffer = io.BytesIO()
    presentation.save(buffer)
    buffer.seek(0)
    return buffer.read()


def _image_to_pdf_bytes(image_bytes: bytes) -> bytes:
    image = Image.open(io.BytesIO(image_bytes))
    if image.mode in ("RGBA", "P"):
        image = image.convert("RGB")

    buffer = io.BytesIO()
    image.save(buffer, format="PDF")
    buffer.seek(0)
    return buffer.read()


def _pdf_to_images_zip(pdf_bytes: bytes, image_format: str) -> bytes:
    fmt = image_format.lower()
    if fmt not in ("jpg", "jpeg", "png"):
        raise ValueError("Format must be jpg or png")

    ext = "jpg" if fmt in ("jpg", "jpeg") else "png"
    pdf = fitz.open(stream=pdf_bytes, filetype="pdf")
    buffer = io.BytesIO()

    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as archive:
        for page_index in range(len(pdf)):
            page = pdf[page_index]
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            image_bytes = pix.tobytes(ext)
            archive.writestr(f"page_{page_index + 1:03d}.{ext}", image_bytes)

    pdf.close()
    buffer.seek(0)
    return buffer.read()


def _pdf_to_text_bytes(pdf_bytes: bytes) -> bytes:
    pdf = fitz.open(stream=pdf_bytes, filetype="pdf")
    parts: list[str] = []

    for page_index in range(len(pdf)):
        page = pdf[page_index]
        text = page.get_text("text").strip()
        if text:
            parts.append(text)
        else:
            parts.append(f"[Page {page_index + 1} — no extractable text]")

    pdf.close()
    return "\n\n".join(parts).encode("utf-8")


def _pdf_to_html_bytes(pdf_bytes: bytes) -> bytes:
    pdf = fitz.open(stream=pdf_bytes, filetype="pdf")
    parts: list[str] = [
        "<!DOCTYPE html>",
        "<html><head><meta charset=\"utf-8\"><title>PDF Export</title></head><body>",
    ]

    for page_index in range(len(pdf)):
        page = pdf[page_index]
        html = page.get_text("html").strip()
        parts.append(f"<section data-page=\"{page_index + 1}\">")
        parts.append(html if html else f"<p>[Page {page_index + 1} — no extractable content]</p>")
        parts.append("</section>")

    parts.append("</body></html>")
    pdf.close()
    return "\n".join(parts).encode("utf-8")


@router.post("/to-pdf")
async def to_pdf(file: UploadFile = File(...)) -> StreamingResponse:
    data, ext = await read_validated_upload(file, TO_PDF_EXTS)

    try:
        if ext in HTML_TO_PDF_EXTS:
            pdf_bytes = html_to_pdf(data)
        else:
            pdf_bytes = convert_to_pdf(data, ext)
    except (OfficeConversionError, HtmlConversionError) as exc:
        raise HTTPException(status_code=503, detail=str(exc)) from exc
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"Conversion failed: {exc}") from exc

    stem = Path(file.filename).stem
    return StreamingResponse(
        io.BytesIO(pdf_bytes),
        media_type=MIME_PDF,
        headers={"Content-Disposition": f'attachment; filename="{stem}.pdf"'},
    )


@router.post("/image-to-pdf")
async def image_to_pdf(file: UploadFile = File(...)) -> StreamingResponse:
    data, _ext = await read_validated_upload(file, IMAGE_TO_PDF_EXTS)

    try:
        pdf_bytes = _image_to_pdf_bytes(data)
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"Conversion failed: {exc}") from exc

    stem = Path(file.filename).stem
    return StreamingResponse(
        io.BytesIO(pdf_bytes),
        media_type=MIME_PDF,
        headers={"Content-Disposition": f'attachment; filename="{stem}.pdf"'},
    )


@router.post("/pdf-to-images")
async def pdf_to_images(
    file: UploadFile = File(...),
    format: str = Form(default="jpg"),
) -> StreamingResponse:
    pdf_bytes = await read_validated_pdf(file)

    try:
        zip_bytes = _pdf_to_images_zip(pdf_bytes, format)
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"Conversion failed: {exc}") from exc

    stem = Path(file.filename).stem
    return StreamingResponse(
        io.BytesIO(zip_bytes),
        media_type=MIME_ZIP,
        headers={"Content-Disposition": f'attachment; filename="{stem}_pages.zip"'},
    )


@router.post("/pdf-to-text")
async def pdf_to_text(file: UploadFile = File(...)) -> StreamingResponse:
    pdf_bytes = await read_validated_pdf(file)

    try:
        text_bytes = _pdf_to_text_bytes(pdf_bytes)
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"Conversion failed: {exc}") from exc

    stem = Path(file.filename).stem
    return StreamingResponse(
        io.BytesIO(text_bytes),
        media_type=MIME_TXT,
        headers={"Content-Disposition": f'attachment; filename="{stem}.txt"'},
    )


@router.post("/pdf-to-html")
async def pdf_to_html(file: UploadFile = File(...)) -> StreamingResponse:
    pdf_bytes = await read_validated_pdf(file)

    try:
        html_bytes = _pdf_to_html_bytes(pdf_bytes)
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"Conversion failed: {exc}") from exc

    stem = Path(file.filename).stem
    return StreamingResponse(
        io.BytesIO(html_bytes),
        media_type=MIME_HTML,
        headers={"Content-Disposition": f'attachment; filename="{stem}.html"'},
    )


@router.post("/pdf-to-word")
async def pdf_to_word(file: UploadFile = File(...)) -> StreamingResponse:
    pdf_bytes = await read_validated_pdf(file)

    try:
        docx_bytes = _pdf_to_docx_bytes(pdf_bytes)
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"Conversion failed: {exc}") from exc

    stem = Path(file.filename).stem
    return StreamingResponse(
        io.BytesIO(docx_bytes),
        media_type=MIME_DOCX,
        headers={"Content-Disposition": f'attachment; filename="{stem}.docx"'},
    )


@router.post("/pdf-to-excel")
async def pdf_to_excel(file: UploadFile = File(...)) -> StreamingResponse:
    pdf_bytes = await read_validated_pdf(file)

    try:
        xlsx_bytes = _pdf_to_xlsx_bytes(pdf_bytes)
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"Conversion failed: {exc}") from exc

    stem = Path(file.filename).stem
    return StreamingResponse(
        io.BytesIO(xlsx_bytes),
        media_type=MIME_XLSX,
        headers={"Content-Disposition": f'attachment; filename="{stem}.xlsx"'},
    )


@router.post("/pdf-to-ppt")
async def pdf_to_ppt(file: UploadFile = File(...)) -> StreamingResponse:
    pdf_bytes = await read_validated_pdf(file)

    try:
        pptx_bytes = _pdf_to_pptx_bytes(pdf_bytes)
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"Conversion failed: {exc}") from exc

    stem = Path(file.filename).stem
    return StreamingResponse(
        io.BytesIO(pptx_bytes),
        media_type=MIME_PPTX,
        headers={"Content-Disposition": f'attachment; filename="{stem}.pptx"'},
    )
